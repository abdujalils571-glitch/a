import logging, asyncio, os, threading, re, sqlite3
from http.server import HTTPServer, BaseHTTPRequestHandler
from groq import Groq
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command, CommandObject
from aiogram.types import (ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup, 
                           InlineKeyboardButton, FSInputFile, CallbackQuery)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- 1. HEALTH CHECK ---
class RenderHandler(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200); self.end_headers()
        self.wfile.write(b"SlideMaster AI Mega Ultra: Active")

def run_health_check():
    server = HTTPServer(('0.0.0.0', int(os.environ.get("PORT", 10000))), RenderHandler)
    server.serve_forever()

# --- 2. CONFIG ---
API_TOKEN = os.environ.get('BOT_TOKEN')
GROQ_API_KEY = os.environ.get('GROQ_API_KEY')
ADMIN_ID = 8049278418 
CHANNEL_ID = "@abdujalils" 

client = Groq(api_key=GROQ_API_KEY)
bot = Bot(token=API_TOKEN)
dp = Dispatcher()
logging.basicConfig(level=logging.INFO)

# --- 3. DATABASE (PRO VERSION) ---
def db_init():
    conn = sqlite3.connect('pptx_master.db')
    conn.execute("""CREATE TABLE IF NOT EXISTS users (
        id BIGINT PRIMARY KEY, lang TEXT, is_premium INT DEFAULT 0, 
        usage_count INT DEFAULT 0, invited_by BIGINT, joined_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )""")
    conn.commit(); conn.close()

def db_query(sql, params=(), commit=False, fetch_all=False):
    try:
        conn = sqlite3.connect('pptx_master.db', timeout=10)
        cursor = conn.cursor()
        # SQL inyeksiyadan himoya va formatlash
        cursor.execute(sql.replace('%s', '?'), params)
        
        if commit: 
            conn.commit()
            return True
        
        # Agar fetch_all True bo'lsa barcha natijalarni, aks holda bittasini qaytaradi
        if fetch_all:
            res = cursor.fetchall()
        else:
            res = cursor.fetchone()
        return res
    except Exception as e:
        logging.error(f"DB Error: {e}")
        return None
    finally: 
        conn.close()
# --- 4. MULTILINGUAL CONTENT (BUSINESS CLASS) ---
LANGS = {
    'uz': {
        'welcome': "✨ **SlideMaster AI Ultra**\n\nProfessional prezentatsiyalar olamiga xush kelibsiz! Mavzuni kiriting, AI qolganini bajaradi.",
        'btns': ["💎 VIP Tariflar", "📊 Kabinet", "🤝 Hamkorlik", "📚 Qo'llanma", "🌐 Tilni o'zgartirish"],
        'sub_err': "🚀 **DIQQAT!**\n\nXizmatdan super-tezkor foydalanish uchun rasmiy kanalimizga a'zo bo'ling.",
        'guide': "📖 **PROFESSIONAL YO'RIQNOMA:**\n\n1️⃣ **Mavzu:** Istalgan tilda mavzu bering.\n2️⃣ **Tanlov:** 10, 15 yoki 20 ta slaydni tanlang.\n3️⃣ **Natija:** AI 100+ manbadan foydalanib slayd tayyorlaydi.\n\n🎁 **Bonus:** Dastlabki 3 ta slayd mutlaqo BEPUL!",
        'tarif': "💎 **PREMIUM PLANLAR:**\n\n🔹 **Standard (1 slayd):** 4,000 UZS\n🔹 **Business (10 slayd):** 15,000 UZS\n🔹 **Ultra (20 slayd):** 25,000 UZS\n🌟 **INFINITY (Cheksiz):** 35,000 UZS\n\n💳 **To'lov:** `9860230107924485` \n👤 **Egasi:** Abdujalil A.\n\n*To'lovdan so'ng chekni adminga yuboring!*",
        'ref': "🚀 **BUSINESS CLASS REFERRAL**\n\n{bar} ({count}/10)\n\nHar bir taklif uchun bonus! 10 do'st = **FREE PREMIUM**\n\n🔗 Sening havolang:\n{link}",
        'wait': "🧠 **AI tahlil qilmoqda...**\nEkspert darajasidagi ma'lumotlar saralanmoqda. Iltimos kuting...",
        'done': "✅ **Tayyor!** Slaydingiz yuqori sifatda tayyorlandi."
    },
    'ru': {
        'welcome': "✨ **SlideMaster AI Ultra**\n\nДобро пожаловать! Введите тему, и ИИ создаст профессиональную презентацию за секунды.",
        'btns': ["💎 VIP Тарифы", "📊 Кабинет", "🤝 Партнерство", "📚 Гайд", "🌐 Язык"],
        'sub_err': "🚀 **ВНИМАНИЕ!**\n\nПодпишитесь на канал для мгновенного доступа к функциям ИИ.",
        'guide': "📖 **ИНСТРУКЦИЯ:**\n\n1️⃣ **Тема:** Любая тема на любом языке.\n2️⃣ **Выбор:** 10, 15 или 20 слайдов.\n3️⃣ **Результат:** Глубокий анализ данных.\n\n🎁 **Бонус:** Первые 3 генерации БЕСПЛАТНО!",
        'tarif': "💎 **ТАРИФНЫЕ ПЛАНЫ:**\n\n🔹 **Standard (1):** 4,000 UZS\n🔹 **Business (10):** 15,000 UZS\n🔹 **Ultra (20):** 25,000 UZS\n🌟 **INFINITY:** 35,000 UZS\n\n💳 **Карта:** `9860230107924485` \n👤 **Владелец:** Абдужалил А.",
        'ref': "🚀 **REFERRAL SYSTEM**\n\n{bar} ({count}/10)\n\n10 друзей = **FREE PREMIUM**\n\n🔗 Ваша ссылка:\n{link}",
        'wait': "🧠 **ИИ анализирует...**\nСобираем лучшие экспертные данные. Пожалуйста, подождите...",
        'done': "✅ **Готово!** Ваша презентация высокого качества готова."
    },
    'en': {
        'welcome': "✨ **SlideMaster AI Ultra**\n\nWelcome to the future! Enter a topic and get a professional presentation instantly.",
        'btns': ["💎 VIP Plans", "📊 Account", "🤝 Referral", "📚 Guide", "🌐 Language"],
        'sub_err': "🚀 **ATTENTION!**\n\nSubscribe to our channel to unlock super-fast AI generation.",
        'guide': "📖 **GUIDE:**\n\n1️⃣ **Topic:** Any topic, any language.\n2️⃣ **Selection:** 10, 15, or 20 slides.\n3️⃣ **Result:** Expert-curated content.\n\n🎁 **Bonus:** First 3 generations are FREE!",
        'tarif': "💎 **PREMIUM PLANS:**\n\n🔹 **Standard (1):** 4,000 UZS\n🔹 **Business (10):** 15,000 UZS\n🔹 **Ultra (20):** 25,000 UZS\n🌟 **INFINITY:** 35,000 UZS\n\n💳 **Card:** `9860230107924485` \n👤 **Owner:** Abdujalil A.",
        'ref': "🚀 **BUSINESS CLASS REFERRAL**\n\n{bar} ({count}/10)\n\n10 friends = **FREE PREMIUM**\n\n🔗 Your link:\n{link}",
        'wait': "🧠 **AI is thinking...**\nCurating expert insights. Please wait...",
        'done': "✅ **Success!** Your high-quality presentation is ready."
    }
}
def get_admin_pay_kb(user_id):
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✅ Standart (+1)", callback_data=f"give_10_{user_id}")],
        [InlineKeyboardButton(text="✅ Business (+10)", callback_data=f"give_15_{user_id}")],
        [InlineKeyboardButton(text="✅ Ultra (+20)", callback_data=f"give_20_{user_id}")],
        [InlineKeyboardButton(text="🌟 INFINITY (Full)", callback_data=f"give_inf_{user_id}")],
        [InlineKeyboardButton(text="❌ Rad etish", callback_data=f"reject_{user_id}")]
    ])

# --- 5. CORE FUNCTIONS (SPEED & QUALITY) ---
async def is_subscribed(uid):
    try:
        m = await bot.get_chat_member(CHANNEL_ID, uid)
        return m.status in ['member', 'administrator', 'creator']
    except: return False

def get_main_kb(lang):
    b = LANGS[lang]['btns']
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=b[0]), KeyboardButton(text=b[1])],
        [KeyboardButton(text=b[2])],
        [KeyboardButton(text=b[3]), KeyboardButton(text=b[4])]
    ], resize_keyboard=True)

def create_ultra_pptx(topic, ai_text, uid):
    try:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        
        # Slaydlarni ajratish (Yaxshilangan regex)
        slides = [s.strip() for s in re.split(r'---|\n(?=Slide|Slayd|Слайд|#)', ai_text) if len(s.strip()) > 20]

        for i, content in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Premium Background (Dark Gradient Style)
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(15, 15, 25)
            bg.line.fill.background()

            lines = [l.strip() for l in content.split("\n") if l.strip()]
            if not lines: continue

            if i == 0: # Title Slide (Ultra Design)
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(3))
                p = title_box.text_frame.paragraphs[0]
                p.text = topic.upper()
                p.font.size, p.font.bold = Pt(54), True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
            else:
                # Header
                header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
                h_p = header_box.text_frame.paragraphs[0]
                h_p.text = lines[0].replace("#", "").strip()[:70]
                h_p.font.size, h_p.font.bold = Pt(34), True
                h_p.font.color.rgb = RGBColor(0, 200, 255) # Cyan Blue

                # Body (Overflow Protection)
                body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))
                tf = body_box.text_frame
                tf.word_wrap = True
                
                for ln in lines[1:8]: # Max 7 bullets to avoid overflow
                    p = tf.add_paragraph()
                    p.text = f"• {ln.lstrip('* -')}"
                    p.font.color.rgb = RGBColor(230, 230, 230)
                    p.font.size = Pt(18) if len(content) < 500 else Pt(15)
                    p.space_before = Pt(8)

        path = f"ultra_{uid}.pptx"; prs.save(path); return path
    except Exception as e:
        logging.error(f"PPTX Error: {e}"); return None

# --- 6. HANDLERS ---
@dp.message(Command("start"))
async def start_cmd(m: types.Message, command: CommandObject):
    uid = m.from_user.id
    u = db_query("SELECT lang FROM users WHERE id=%s", (uid,))
    
    if not u:
        ref_id = int(command.args) if command.args and command.args.isdigit() else None
        db_query("INSERT INTO users (id, lang, invited_by) VALUES (%s, 'uz', %s)", (uid, ref_id), commit=True)
        u = ('uz',)
        if ref_id:
            await bot.send_message(ref_id, "🤝 **Yangi do'st taklif qilindi!**")

    if not await is_subscribed(uid):
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📢 Kanalga a'zo bo'lish", url=f"https://t.me/{CHANNEL_ID[1:]}")],
            [InlineKeyboardButton(text="✅ Tekshirish", callback_data="check")]
        ])
        return await m.answer(LANGS[u[0]]['sub_err'], reply_markup=kb)
    
    await m.answer(LANGS[u[0]]['welcome'], reply_markup=get_main_kb(u[0]))

@dp.message(F.photo)
async def handle_payment_check(m: types.Message):
    uid = m.from_user.id
    # Foydalanuvchiga javob
    await m.reply("⏳ **Chek qabul qilindi!**\nAdmin tasdiqlashini kuting. Tez orada javob beramiz.")
    
    # Adminga yuborish
    caption = (f"💰 **Yangi to'lov cheki!**\n\n"
               f"👤 Foydalanuvchi: {m.from_user.full_name}\n"
               f"🆔 ID: `{uid}`\n"
               f"🔗 Link: tg://user?id={uid}")
    
    await bot.send_photo(
        chat_id=ADMIN_ID, 
        photo=m.photo[-1].file_id, 
        caption=caption, 
        reply_markup=get_admin_pay_kb(uid)
    )
@dp.message(F.text)
async def handle_text(m: types.Message):
    uid = m.from_user.id
    res = db_query("SELECT lang, usage_count, is_premium FROM users WHERE id=%s", (uid,))
    if not res: return
    l = res[0]

    # Admin Panel
    if m.text == "/admin" and uid == ADMIN_ID:
        total = db_query("SELECT COUNT(*) FROM users")[0]
        return await m.answer(f"📊 **ADMIN PANEL**\n\nJami foydalanuvchilar: {total}\n\nReklama uchun: `/send [xabar]`")

    if m.text.startswith("/send") and uid == ADMIN_ID:
        # Agar admin biror xabarga reply qilib /send yozsa, o'sha xabarni yuboradi
        if not m.reply_to_message:
            return await m.answer("❗ **Xatolik:** Reklama yuborish uchun biror xabarga (rasm, matn, video) `/send` deb reply qiling.")

        users = db_query("SELECT id FROM users", fetch_all=True)
        if not users:
            return await m.answer("Foydalanuvchilar topilmadi.")

        status_msg = await m.answer(f"🚀 **Reklama yuborish boshlandi...**\nJami: {len(users)} ta manzil.")
        
        done = 0
        blocked = 0
        errors = 0

        for user in users:
            try:
                # copy_message xabarni asl holatida (rasm, tugma, caption) nusxalaydi
                await bot.copy_message(
                    chat_id=user[0],
                    from_chat_id=m.chat.id,
                    message_id=m.reply_to_message.message_id
                )
                done += 1
                
                # Har 30 ta xabarda adminni xabardor qilish
                if done % 30 == 0:
                    await status_msg.edit_text(f"⏳ **Yuborilmoqda...**\n\n✅ Yetkazildi: {done}\n🚫 Bloklangan: {blocked}")
                
                # Telegram cheklovlaridan qochish uchun kichik tanaffus
                await asyncio.sleep(0.05) 
                
            except Exception as e:
                if "bot was blocked" in str(e).lower():
                    blocked += 1
                else:
                    errors += 1
                    logging.error(f"Send Error for {user[0]}: {e}")

        final_text = (f"✅ **Reklama yakunlandi!**\n\n"
                      f"📊 **Statistika:**\n"
                      f"🟢 Yetkazildi: {done}\n"
                      f"🔴 Bloklangan (Botni o'chirgan): {blocked}\n"
                      f"⚠️ Boshqa xatoliklar: {errors}\n\n"
                      f"🏁 Jami: {len(users)}")
        
        return await status_msg.edit_text(final_text)
    
    if m.text in [LANGS['uz']['btns'][0], LANGS['ru']['btns'][0], LANGS['en']['btns'][0]]:
        await m.answer(LANGS[l]['tarif'], parse_mode="Markdown")
    
    elif m.text in [LANGS['uz']['btns'][1], LANGS['ru']['btns'][1], LANGS['en']['btns'][1]]:
        cnt = db_query("SELECT COUNT(*) FROM users WHERE invited_by=%s", (uid,))[0]
        await m.answer(f"👤 **KABINET**\n\n🆔 ID: `{uid}`\n💎 Status: {'Premium' if res[2] else 'Bepul'}\n📊 Foydalanildi: {res[1]}\n👥 Do'stlar: {cnt}")

    elif m.text in [LANGS['uz']['btns'][2], LANGS['ru']['btns'][2], LANGS['en']['btns'][2]]:
        cnt = db_query("SELECT COUNT(*) FROM users WHERE invited_by=%s", (uid,))[0]
        link = f"https://t.me/{(await bot.get_me()).username}?start={uid}"
        bar = ("🔵" * min(cnt, 10)) + ("⚪" * (10 - min(cnt, 10)))
        await m.answer(LANGS[l]['ref'].format(bar=bar, count=cnt, link=link))

    elif m.text in [LANGS['uz']['btns'][3], LANGS['ru']['btns'][3], LANGS['en']['btns'][3]]:
        await m.answer(LANGS[l]['guide'], parse_mode="Markdown")

    elif m.text in [LANGS['uz']['btns'][4], LANGS['ru']['btns'][4], LANGS['en']['btns'][4]]:
        kb = InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="🇺🇿 UZ", callback_data="setlang_uz"),
            InlineKeyboardButton(text="🇷🇺 RU", callback_data="setlang_ru"),
            InlineKeyboardButton(text="🇺🇸 EN", callback_data="setlang_en")
        ]])
        await m.answer("Select language / Tilni tanlang:", reply_markup=kb)

    elif not m.text.startswith("/"):
        # Slayd miqdorini tanlash
        kb = InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="10 Slayd", callback_data=f"gen_10_{m.text[:25]}"),
            InlineKeyboardButton(text="15 Slayd", callback_data=f"gen_15_{m.text[:25]}"),
            InlineKeyboardButton(text="20 Slayd (VIP)", callback_data=f"gen_20_{m.text[:25]}")
        ]])
        await m.answer(f"🎯 **Mavzu:** {m.text}\nSlaydlar sonini tanlang:", reply_markup=kb)
@dp.callback_query(F.data.startswith(("give_", "reject_")))
async def process_admin_decision(c: CallbackQuery):
    if c.from_user.id != ADMIN_ID: return
    
    data = c.data.split("_")
    action = data[1] # 10, 15, 20, inf, reject
    target_id = int(data[2])

    if action == "reject":
        await bot.send_message(target_id, "❌ **To'lovingiz tasdiqlanmadi.**\nIltimos, chekni qayta tekshirib yuboring yoki adminga murojaat qiling.")
        return await c.message.edit_caption(caption="🚫 Rad etildi")

    # Tariflar bo'yicha status berish
    if action == "inf":
        db_query("UPDATE users SET is_premium=1, usage_count=0 WHERE id=%s", (target_id,), commit=True)
        msg = "🌟 Tabriklaymiz! Sizga **INFINITY (Cheksiz)** paketi berildi!"
    else:
        plus_count = int(action)
        # Hozirgi limitiga qo'shish (limitni kamaytirish orqali premium berish)
        db_query("UPDATE users SET usage_count = usage_count - %s WHERE id=%s", (plus_count, target_id), commit=True)
        msg = f"✅ Tabriklaymiz! Sizning balansingizga **{plus_count} ta slayd** qo'shildi!"

    await bot.send_message(target_id, msg)
    await c.message.edit_caption(caption=f"✅ Tasdiqlandi: {action}")

@dp.callback_query(F.data.startswith("gen_"))
async def generate_callback(c: CallbackQuery):
    _, count, topic = c.data.split("_", 2)
    uid = c.from_user.id
    u = db_query("SELECT usage_count, is_premium, lang FROM users WHERE id=%s", (uid,))
    
    # FREE CHANCE: Dastlabki 3 ta tekin
    if not u[1] and u[0] >= 3:
        return await c.answer("⚠️ Bepul imkoniyat tugadi. Premiumga o'ting!", show_alert=True)
    
    wait_msg = await c.message.edit_text(LANGS[u[2]]['wait'])
    
    try:
        # SUPER AI PROMPT
        prompt = (f"Act as a World-Class Consultant. Create a detailed presentation for: {topic}. "
                  f"Language: {u[2]}. Total slides: {count}. "
                  "Structure: Use '# Title' for slide headers and bullet points for content. "
                  "Each slide must have 5-7 high-quality, professional insights. "
                  "Use '---' to separate slides. Do not include any intro/outro text.")
        
        response = await asyncio.to_thread(client.chat.completions.create, 
                                          model="llama-3.3-70b-versatile", 
                                          messages=[{"role":"system","content":prompt}])
        
        path = await asyncio.to_thread(create_ultra_pptx, topic, response.choices[0].message.content, uid)
        
        if path:
            await bot.send_document(uid, FSInputFile(path), caption=LANGS[u[2]]['done'])
            db_query("UPDATE users SET usage_count=usage_count+1 WHERE id=%s", (uid,), commit=True)
            os.remove(path)
        else:
            await c.message.answer("❌ Error creating PPTX")
    except Exception as e:
        await c.message.answer(f"❌ AI Error: {e}")
    finally:
        await wait_msg.delete()

@dp.callback_query(F.data.startswith("setlang_"))
async def setlang(c: CallbackQuery):
    l = c.data.split("_")[1]
    db_query("UPDATE users SET lang=%s WHERE id=%s", (l, c.from_user.id), commit=True)
    await c.message.delete()
    await bot.send_message(c.from_user.id, "✅ Done / Tayyor!", reply_markup=get_main_kb(l))

@dp.callback_query(F.data == "check")
async def check_subscription(c: CallbackQuery):
    if await is_subscribed(c.from_user.id):
        await c.answer("✅ Rahmat!"); await c.message.delete()
        u = db_query("SELECT lang FROM users WHERE id=%s", (c.from_user.id,))
        await bot.send_message(c.from_user.id, LANGS[u[0]]['welcome'], reply_markup=get_main_kb(u[0]))
    else:
        await c.answer("❌ Hali a'zo emassiz!", show_alert=True)

# --- 7. STARTUP ---
async def main():
    db_init()
    threading.Thread(target=run_health_check, daemon=True).start()
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main()) KODNI BIROR JOYIDA XATO BORMI